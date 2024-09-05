import os
import streamlit as st
import PyPDF2
import docx
from io import StringIO
from pptx import Presentation
from pptx.util import Inches
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain

# Set up Google API Key directly in the code
GOOGLE_API_KEY = "AIzaSyBoX4UUHV5FO4lvYwdkSz6R5nlxLadTHnU"

# Initialize the Gemini 1.5 Flash model
def init_llm(api_key):
    return ChatGoogleGenerativeAI(
        model="gemini-1.5-flash",  # Using Gemini 1.5 Flash as per request
        api_key=api_key,
        temperature=0.8,  # Adjust temperature for creative responses
        max_tokens=300,    # Adjust max tokens for better performance
        timeout=15,        # Timeout after 15 seconds
        max_retries=3      # Retry up to 3 times in case of errors
    )

# Helper function to extract text from PDFs
def extract_text_from_pdf(pdf_file):
    reader = PyPDF2.PdfReader(pdf_file)
    text = ""
    for page_num in range(len(reader.pages)):
        text += reader.pages[page_num].extract_text()
    return text

# Helper function to extract text from text files or docx files
def extract_text_from_file(file):
    if file.type == "text/plain":
        return StringIO(file.getvalue().decode("utf-8")).read()
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = docx.Document(file)
        return "\n".join([para.text for para in doc.paragraphs])
    return ""

# Function to create a PowerPoint file from AI-generated content
def create_ppt_from_text(text):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Use a layout with a title and content

    # Split the text into sections based on some logic
    sections = text.split("\n\n")  # Basic splitting by paragraphs for now

    for i, section in enumerate(sections):
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"Slide {i + 1}"
        content = slide.shapes.placeholders[1]
        content.text = section.strip()  # Add the section content

    # Save PowerPoint
    ppt_filename = "generated_presentation.pptx"
    prs.save(ppt_filename)
    return ppt_filename

# Streamlit interface
st.title("Document Processing App")

# File upload or manual text input
uploaded_file = st.file_uploader("Upload a PDF or text file", type=["pdf", "txt", "docx"])
manual_text = st.text_area("Or manually enter the text here")

if uploaded_file:
    if uploaded_file.type == "application/pdf":
        document_text = extract_text_from_pdf(uploaded_file)
    else:
        document_text = extract_text_from_file(uploaded_file)
elif manual_text:
    document_text = manual_text
else:
    st.write("Please upload a file or enter text manually.")
    st.stop()

# Let the user choose the action
option = st.selectbox("Choose what you want to do with the text:", 
                      ["", "Create Summary", "Q&A Chatbot", "Create PPT from Text"])

# Initialize LLM
llm = init_llm(GOOGLE_API_KEY)

# Handle the user's choice
if option == "Create Summary":
    with st.spinner("Generating summary..."):
        summary_prompt = f"Summarize the following text: {document_text}"
        response = llm.invoke(summary_prompt)
        st.write("### Summary:")
        st.write(response.content)

elif option == "Q&A Chatbot":
    st.write("Ask questions about the document.")
    question = st.text_input("Enter your question:")
    
    if st.button("Ask"):
        with st.spinner("Fetching answer..."):
            qa_prompt = f"Based on this document: {document_text}\nAnswer the question: {question}"
            response = llm.invoke(qa_prompt)
            st.write("### Answer:")
            st.write(response.content)

elif option == "Create PPT from Text":
    # Ask the user for the min and max slides once they select Create PPT
    min_slides = st.number_input("Enter the minimum number of slides", min_value=1, max_value=10, value=4)
    max_slides = st.number_input("Enter the maximum number of slides", min_value=min_slides, max_value=20, value=6)

    if st.button("Create PPT"):
        with st.spinner("Creating detailed PowerPoint..."):
            # Optimized prompt for the model
            ppt_prompt = (
                f"Create a visually appealing PowerPoint presentation from the following document. "
                f"The presentation should be between {min_slides} and {max_slides} slides long:\n\n{document_text}"
            )
            response = llm.invoke(ppt_prompt)
            
            # Generate PowerPoint file using AI-generated content
            ppt_file = create_ppt_from_text(response.content)
            
            with open(ppt_file, "rb") as f:
                st.download_button("Download PPT", f, file_name="generated_presentation.pptx")
